"""
Build IonShield investor pitch deck (.pptx).

Aesthetic: Palantir/Anduril — dark background, white type, single cyan accent,
no emoji, no animations, generous whitespace, serif-free.

Run:
    python3 scripts/build_pitch_deck.py
Output:
    IonShield_PitchDeck.pptx in repo root
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0E, 0x14)   # near-black
PANEL    = RGBColor(0x12, 0x18, 0x20)
RULE     = RGBColor(0x1F, 0x2A, 0x36)
TEXT     = RGBColor(0xE6, 0xEC, 0xF2)
MUTED    = RGBColor(0x8B, 0x97, 0xA6)
ACCENT   = RGBColor(0x4E, 0xC9, 0xE6)   # cyan
WARN     = RGBColor(0xE8, 0xB0, 0x4A)   # amber, used sparingly

FONT = "Helvetica Neue"

# ── presentation setup ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, items, left, top, width, height,
                size=16, color=TEXT, line_spacing=1.25, bullet_char="—"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = f"{bullet_char}  {item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_rule(slide, left, top, width, color=ACCENT, weight=2.0):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(weight * 9525)))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = color
    return line


def add_panel(slide, left, top, width, height, color=PANEL):
    p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    p.line.color.rgb = RULE
    p.line.width = Pt(0.75)
    p.fill.solid()
    p.fill.fore_color.rgb = color
    return p


def footer(slide, pageno, total=15):
    add_text(slide, "IonShield", Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
             size=9, color=MUTED)
    add_text(slide, f"{pageno} / {total}", Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def std_slide(title, eyebrow=None):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    if eyebrow:
        add_text(s, eyebrow.upper(), Inches(0.6), Inches(0.45), Inches(10), Inches(0.3),
                 size=10, bold=True, color=ACCENT)
    add_text(s, title, Inches(0.6), Inches(0.75), Inches(12), Inches(0.7),
             size=30, bold=True, color=TEXT)
    add_rule(s, Inches(0.6), Inches(1.45), Inches(1.2))
    return s


# ── Slide 1: Cover ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
# subtle accent bar
add_rule(s, Inches(0.6), Inches(0.6), Inches(0.6), color=ACCENT, weight=3.0)
add_text(s, "IonShield", Inches(0.6), Inches(2.6), Inches(12), Inches(1.4),
         size=72, bold=True, color=TEXT)
add_text(s,
         "Operational intelligence for contested electromagnetic environments.",
         Inches(0.6), Inches(4.0), Inches(12), Inches(0.6),
         size=22, color=MUTED)
add_text(s,
         "Pre-seed  ·  Founded 2025  ·  Jacob West, Founder",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         size=12, color=MUTED)
add_notes(s,
          "IonShield turns raw space-weather and ionospheric data into mission-grade "
          "decisions for operators flying drones, planning patrols, running RTK ag, "
          "or coordinating contested-EM operations. We are raising pre-seed.")

# ── Slide 2: Problem ─────────────────────────────────────────────────────────
s = std_slide("GPS, HF, and SATCOM degrade silently — and operators find out mid-mission.",
              eyebrow="Problem")
add_bullets(s, [
    "Modern operations assume clean GPS, reliable HF, and stable SATCOM. None of that is guaranteed.",
    "Space-weather and ionospheric disturbances degrade RF systems daily — most operators have no visibility into it.",
    "Existing space-weather data exists, but it is scientific. Kp = 6 does not tell a UAV operator whether to fly.",
    "Defense, autonomy, ag, maritime, and aviation all hit the same wall: raw data, no decision.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=18)
footer(s, 2)
add_notes(s,
          "Space weather is not a niche concern any more. Autonomous systems, BVLOS UAVs, "
          "RTK precision ag, and modern defense workflows all assume GPS, HF, or SATCOM. "
          "All of them degrade in ways the operator cannot see until something goes wrong.")

# ── Slide 3: Why Now ─────────────────────────────────────────────────────────
s = std_slide("Three forces converging in the next 24 months.", eyebrow="Why now")
add_bullets(s, [
    "Solar Cycle 25 peak — geomagnetic storms at multi-decade highs through 2026–2027.",
    "Autonomy at scale — BVLOS drones, autonomous ground, maritime, and ag are deploying now, all GPS-dependent.",
    "Contested EM as doctrine — GPS jamming and HF disruption are now standard adversary capability, not edge cases.",
    "Operators need a translation layer between physical reality and mission decisions. Today, there isn't one.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=18)
footer(s, 3)
add_notes(s,
          "Solar max, autonomy at scale, and contested EM as standard doctrine all hit at once. "
          "Each on its own would create demand. Together they create a category.")

# ── Slide 4: Why existing solutions fail operators ──────────────────────────
s = std_slide("Existing tools are scientific, not operational.", eyebrow="Gap")
# 2-column compare
left_x = Inches(0.6); right_x = Inches(6.9); col_w = Inches(6.0); col_y = Inches(1.9)
add_panel(s, left_x, col_y, col_w, Inches(4.6))
add_panel(s, right_x, col_y, col_w, Inches(4.6))
add_text(s, "What exists today", left_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_bullets(s, [
    "NOAA SWPC — authoritative, but raw scientific indices (Kp, Bz, X-ray flux).",
    "Tomorrow.io space-weather — alert feeds, not mission decisions.",
    "Mission Space — premium forecasting for satellite operators, not ground/air operators.",
    "Internal SA tools (defense) — siloed, not integrated with autonomy stacks.",
], left_x + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.4), Inches(3.8), size=13, color=TEXT)

add_text(s, "What an operator actually needs", right_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_bullets(s, [
    "Will my GPS hold accuracy on this route, at this time?",
    "Is my HF link viable for the next 4 hours?",
    "Should I delay, reroute, or proceed?",
    "Can I drop this into ATAK, my GCS, or my fleet planner — right now?",
], right_x + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.4), Inches(3.8), size=13, color=TEXT)
footer(s, 4)
add_notes(s,
          "The data exists. The translation layer does not. Operators are forced to be amateur "
          "space physicists. We close that gap.")

# ── Slide 5: Product overview ────────────────────────────────────────────────
s = std_slide("A decision layer between physics and the mission.", eyebrow="Product")
add_text(s,
         "IonShield ingests authoritative space-weather feeds, runs physics-based degradation "
         "models, and outputs operator-grade decisions: GPS error in metres, HF blackout in "
         "minutes, route verdicts in plain language.",
         Inches(0.6), Inches(1.9), Inches(12), Inches(1.2), size=16, color=TEXT)

# three-pillar row
y = Inches(3.4); w = Inches(4.0); h = Inches(3.5); gap = Inches(0.15)
xs = [Inches(0.6), Inches(0.6) + w + gap, Inches(0.6) + 2 * (w + gap)]
pillars = [
    ("Ingest",
     ["NOAA SWPC, NASA OMNI, GFZ Potsdam",
      "5-minute cadence",
      "SHA-256 provenance on every input"]),
    ("Model",
     ["Klobuchar+Mannucci GPS error",
      "CCIR-888 HF, Bailey PCA, Nakagami-m SATCOM",
      "Pure-Python ML champion/challenger"]),
    ("Decide",
     ["Mission verdicts: CLEAR / CAUTION / HIGH RISK / DELAY",
      "GNSS reliability score, comms risk score",
      "Outputs: ATAK CoT, KML, KMZ, Parquet, JSON"]),
]
for x, (h1, items) in zip(xs, pillars):
    add_panel(s, x, y, w, h)
    add_text(s, h1, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
             size=18, bold=True, color=ACCENT)
    add_bullets(s, items, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), h - Inches(1.2),
                size=13, color=TEXT)
footer(s, 5)
add_notes(s,
          "Three layers: ingest authoritative feeds, run physics-based degradation models, "
          "emit operator decisions. Everything is provenance-tagged.")

# ── Slide 6: Workflow integration (ATAK + beyond) ───────────────────────────
s = std_slide("We meet operators inside the tools they already use.", eyebrow="Integration")
add_bullets(s, [
    "ATAK / CoT 2.0 — IonShield publishes risk overlays directly into Team Awareness Kit.",
    "KML / KMZ — drops into ArcGIS, QGIS, fleet planners, mission planning suites.",
    "Parquet / JSON — data-science and analytics pipelines.",
    "REST API + bearer-token tenants — embed inside autonomy stacks, GCS, and fleet ops.",
    "Foundry-compatible ontology pack — for primes and defense integrators.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 6)
add_notes(s,
          "ATAK-first is deliberate. It is the de facto situational-awareness layer for US "
          "defense and a growing set of dual-use operators. KML/KMZ covers civilian fleet "
          "and surveying workflows. The API covers everything else.")

# ── Slide 7: Defensibility / feedback loop ──────────────────────────────────
s = std_slide("Every customer makes the next decision better.", eyebrow="Defensibility")
add_bullets(s, [
    "Mission-aware thresholds — RTK 0.5 m, high-precision 5 m, standard 10 m, permissive 25 m.",
    "Customer-specific tuning — per-tenant tolerance profiles, per-asset GNSS dependence.",
    "Audit log on every decision — input hash, model version, source labels (measured / modeled / heuristic).",
    "Champion/challenger ML — every new model is A/B-validated against the production champion before promotion.",
    "Compounding moat: feedback loop + provenance + integration depth, not just a model.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 7)
add_notes(s,
          "The moat is not the physics model — those are public. The moat is the translation "
          "layer, the tolerance profiles per mission, the integration depth, and the audit "
          "trail. Each customer sharpens the next decision.")

# ── Slide 8: Market opportunity (dual-use) ──────────────────────────────────
s = std_slide("Dual-use — defense pulls, commercial scales.", eyebrow="Market")
# nine verticals as panels
verticals = [
    ("Defense / NatSec",        "Contested EM, ATAK overlays, mission planning"),
    ("UAV / BVLOS",             "GPS reliability for autonomous flight"),
    ("Precision agriculture",   "RTK tolerance, planting / spraying windows"),
    ("Maritime",                "HF / SATCOM viability, polar routing"),
    ("Aviation",                "Polar and trans-polar route advisories"),
    ("Autonomous ground",       "GPS-dependent fleet ops"),
    ("Surveying / geomatics",   "Centimetre-grade RTK windows"),
    ("Satellite operators",     "Drag, charging, comms degradation"),
    ("Insurance / reinsurance", "Outage modelling, parametric triggers"),
]
cols, rows = 3, 3
gx, gy = Inches(0.6), Inches(1.9)
cw, ch = Inches(4.0), Inches(1.55)
gap_x, gap_y = Inches(0.15), Inches(0.15)
for i, (name, sub) in enumerate(verticals):
    r, c = divmod(i, cols)
    x = gx + c * (cw + gap_x)
    y = gy + r * (ch + gap_y)
    add_panel(s, x, y, cw, ch)
    add_text(s, name, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=ACCENT)
    add_text(s, sub, x + Inches(0.25), y + Inches(0.7), cw - Inches(0.5), Inches(0.8),
             size=11, color=MUTED)
footer(s, 8)
add_notes(s,
          "Multi-vertical by design. Defense is the wedge — willingness to pay, urgency, and "
          "validation. Commercial verticals (UAV, ag, maritime) are the scale path. The "
          "translation layer is the same product for all of them.")

# ── Slide 9: Competitive landscape ──────────────────────────────────────────
s = std_slide("Adjacent — but no one ships operator decisions.", eyebrow="Competition")
# table-style panel
y = Inches(1.9); x = Inches(0.6); w = Inches(12.1); rh = Inches(0.55)
add_panel(s, x, y, w, rh, color=PANEL)
headers = ["", "Authoritative data", "Mission decisions", "ATAK / CoT", "Multi-vertical"]
col_widths = [Inches(3.5), Inches(2.15), Inches(2.15), Inches(2.15), Inches(2.15)]
cx = x
for hh, cw in zip(headers, col_widths):
    add_text(s, hh, cx + Inches(0.15), y + Inches(0.12), cw - Inches(0.2), Inches(0.4),
             size=12, bold=True, color=ACCENT)
    cx += cw
rows = [
    ("NOAA SWPC",            "Yes", "No",      "No",      "n/a"),
    ("Tomorrow.io",          "Yes", "Alerts",  "No",      "Weather-led"),
    ("Mission Space",        "Yes", "Sat ops", "No",      "Satellite"),
    ("Defense SA tools",     "Partial", "Siloed", "Some", "Defense only"),
    ("IonShield",            "Yes", "Yes",     "Yes",     "Yes"),
]
for i, row in enumerate(rows):
    ry = y + rh * (i + 1)
    add_panel(s, x, ry, w, rh, color=BG if i % 2 == 0 else PANEL)
    cx = x
    for j, val in enumerate(row):
        is_us = row[0] == "IonShield"
        color = ACCENT if is_us else TEXT
        bold = is_us
        add_text(s, val, cx + Inches(0.15), ry + Inches(0.12), col_widths[j] - Inches(0.2), Inches(0.4),
                 size=12, bold=bold, color=color)
        cx += col_widths[j]
footer(s, 9)
add_notes(s,
          "Mission Space is the closest analogue, but they sell to satellite operators in "
          "orbit. We sell to operators on the ground and in the air. NOAA is upstream of "
          "everyone. Tomorrow.io is a weather company. No one ships operator decisions today.")

# ── Slide 10: Business model ────────────────────────────────────────────────
s = std_slide("SaaS + integration revenue. Defense anchors, commercial scales.", eyebrow="Business model")
add_bullets(s, [
    "Tiered SaaS per tenant — Operator, Team, Enterprise. Per-seat + per-API-call.",
    "Defense pilots — paid pilots, then ATO-path enterprise contracts. Six-figure ACV.",
    "Commercial — UAV / ag / maritime fleet subscriptions. Four-to-five-figure ACV, high volume.",
    "Integration revenue — Foundry ontology pack, ATAK plugin, GCS partners.",
    "Margins — software margins; physics models are cheap to run, data feeds are public.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 10)
add_notes(s,
          "Defense is the anchor revenue: high ACV, slower sales cycle, but validation gold. "
          "Commercial is the scale layer: lower ACV, faster cycle, much higher volume. "
          "Integration revenue is the third leg.")

# ── Slide 11: Go-to-market ──────────────────────────────────────────────────
s = std_slide("Land via defense pilots, expand into adjacent operators.", eyebrow="Go-to-market")
add_bullets(s, [
    "Phase 1 — paid defense pilots via WarHacker / FI Defense / DRF networks. ATAK-first.",
    "Phase 2 — UAV / BVLOS operators (Part 107 + waiver holders) and RTK ag co-ops.",
    "Phase 3 — primes and integrators (Foundry pack, ATAK partner ecosystem).",
    "Channel: defense accelerators, autonomy-platform partnerships, surveying associations.",
    "Sales motion: founder-led through Series A, then defense BD hire + commercial CSM.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 11)
add_notes(s,
          "Defense pilots first because they validate, pay, and unlock the right networks. "
          "Commercial operators follow because the same translation layer answers their "
          "questions too.")

# ── Slide 12: Validation & traction ─────────────────────────────────────────
s = std_slide("Early signals from the right rooms.", eyebrow="Validation & traction")
add_bullets(s, [
    "Defense Unicorns WarHacker — selected, June 2026.",
    "Founders Institute — Defense & National Security fellowship, Fall 2026.",
    "Dorm Room Fund — pre-seed conversations engaged, pending pilot signal.",
    "Product — full v3 platform live: ingestion, decision engine, mission planner, dashboard, ATAK overlays.",
    "Honest read: pre-revenue, pre-pilot. Next 6 months are about converting these rooms into paid pilots.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 12)
add_notes(s,
          "This is the weakest slide today and we know it. We are pre-revenue and pre-pilot. "
          "What we have is the platform built, the right networks engaged, and a clear plan "
          "to convert these into paid pilots in the next 6 months.")

# ── Slide 13: Founder ───────────────────────────────────────────────────────
s = std_slide("Built by an undergrad physicist working national-security problems.", eyebrow="Founder")
add_text(s, "Jacob West — Founder", Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         size=20, bold=True, color=TEXT)
add_bullets(s, [
    "Undergraduate physics — ionospheric modelling and RF propagation focus.",
    "National-security oriented — selected for Defense Unicorns WarHacker and FI Defense & National Security.",
    "Built the IonShield platform end-to-end: ingestion, physics models, ML, API, dashboard, ATAK integration.",
    "Engineering philosophy: provenance on everything, honest uncertainty, no over-claiming.",
], Inches(0.6), Inches(2.6), Inches(12), Inches(3.8), size=16)
footer(s, 13)
add_notes(s,
          "Solo technical founder today. Hiring plan post-seed: defense BD, ML engineer, "
          "design partner-facing applied engineer.")

# ── Slide 14: Vision & roadmap ──────────────────────────────────────────────
s = std_slide("From decision layer to operational nervous system for contested EM.", eyebrow="Vision")
add_bullets(s, [
    "Now — physics + ML decision engine, mission planner, ATAK / KML / KMZ / Parquet outputs.",
    "6 months — paid defense pilots, first commercial design partners, ATAK plugin GA.",
    "12 months — Foundry pack with primes, RTK-ag and BVLOS commercial GTM, ML champion-2.",
    "24 months — operational nervous system: contested EM, GPS jamming, HF disruption, SATCOM, all in one decision layer.",
], Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), size=17)
footer(s, 14)
add_notes(s,
          "We start with space weather because it is the wedge no one is serving. We expand "
          "into adjacent contested-EM problems — jamming, HF disruption, SATCOM — because "
          "the operator question is the same: can I run this mission, right now?")

# ── Slide 15: The Ask ───────────────────────────────────────────────────────
s = std_slide("Pre-seed — to convert validation into paid pilots.", eyebrow="The ask")
# left: ask figures, right: use of funds
left_x, right_x = Inches(0.6), Inches(7.0)
col_w = Inches(5.7); col_y = Inches(1.9); col_h = Inches(4.6)
add_panel(s, left_x, col_y, col_w, col_h)
add_panel(s, right_x, col_y, col_w, col_h)

add_text(s, "Round", left_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_bullets(s, [
    "Raising — pre-seed",
    "Use of proceeds — 12 months of runway to first paid pilots",
    "Lead — open; introductions welcomed",
    "Vehicle — SAFE",
], left_x + Inches(0.3), col_y + Inches(0.75), col_w - Inches(0.4), Inches(3.5),
   size=14, color=TEXT)

add_text(s, "Use of funds", right_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_bullets(s, [
    "Defense BD lead — pilot conversion",
    "Applied ML engineer — challenger models, per-tenant tuning",
    "Design-partner engineer — ATAK plugin GA, integration depth",
    "Pilot infrastructure — tenant isolation, audit, ATO-path readiness",
], right_x + Inches(0.3), col_y + Inches(0.75), col_w - Inches(0.4), Inches(3.5),
   size=14, color=TEXT)

add_text(s, "jacob@ionshield  ·  ionshield platform live  ·  pilots in flight Q3",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         size=12, color=MUTED, align=PP_ALIGN.LEFT)
footer(s, 15)
add_notes(s,
          "We are raising a pre-seed round to convert the WarHacker / FI Defense / DRF "
          "networks into paid pilots, harden the platform for tenant deployments, and "
          "ship the ATAK plugin to GA. Honest, capital-efficient round.")

# ── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "IonShield_PitchDeck.pptx"
prs.save(out)
print(f"Wrote {out}")
